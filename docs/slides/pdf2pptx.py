"""
Convert presentation.pdf to presentation.pptx
Each PDF page becomes a full-slide image in PowerPoint.

Requirements:
    pip install pdf2image python-pptx Pillow

On Windows you also need poppler:
    1. Download from https://github.com/osber/poppler-windows/releases
    2. Extract and add the 'bin' folder to PATH
    OR use: conda install -c conda-forge poppler
"""

import sys
from pathlib import Path

try:
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Emu, Inches
except ImportError:
    print("Missing dependencies. Install with:")
    print("  pip install pdf2image python-pptx Pillow")
    sys.exit(1)


def pdf_to_pptx(pdf_path: str, output_path: str, dpi: int = 1200):
    pdf_path = Path(pdf_path)
    output_path = Path(output_path)

    print(f"Converting {pdf_path.name} at {dpi} DPI...")
    images = convert_from_path(str(pdf_path), dpi=dpi)
    print(f"  {len(images)} pages extracted.")

    # Create 16:9 presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for i, img in enumerate(images, 1):
        slide = prs.slides.add_slide(blank_layout)

        # Save page as temporary PNG
        tmp_img = pdf_path.parent / f"_tmp_slide_{i}.png"
        img.save(str(tmp_img), "PNG")

        # Add image filling the entire slide
        slide.shapes.add_picture(
            str(tmp_img),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        # Clean up temp file
        tmp_img.unlink()

        if i % 5 == 0:
            print(f"  Processed {i}/{len(images)} slides...")

    prs.save(str(output_path))
    print(f"Done! Saved to: {output_path}")


if __name__ == "__main__":
    script_dir = Path(__file__).parent
    pdf_file = script_dir / "slides.pdf"
    pptx_file = script_dir / "slides.pptx"

    if not pdf_file.exists():
        print(
            f"Error: {pdf_file} not found. Compile the LaTeX first."
        )
        sys.exit(1)

    pdf_to_pptx(str(pdf_file), str(pptx_file))
